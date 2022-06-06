import sys
from PyQt5.QtWidgets import *
from PyQt5 import uic
from pptx import Presentation
import re
from PyQt5.QtGui import QIcon

form_class = uic.loadUiType("design.ui")[0]

class WindowClass(QMainWindow, form_class) :
    def __init__(self) :
        super().__init__()
        self.setupUi(self)
        self.initUI()

        #버튼에 기능 연결
        self.pushButton.clicked.connect(self.makePPT)
        self.comboBox.currentIndexChanged.connect(self.comboBoxTextChange)
        self.comboBox_3.currentIndexChanged.connect(self.comboBox_3TextChange)

        #초기 세팅
        for i in range(50):
            self.comboBox_3.addItem(str(i+1))

    def initUI(self):
        self.setWindowTitle('Bible2PPT')
        self.setWindowIcon(QIcon('icon.png'))

    def comboBoxTextChange(self):
        self.comboBox_3.clear()
        self.comboBox_4.clear()
        self.comboBox_6.clear()
        
        num = self.comboBox.currentIndex()
        if self.comboBox.currentText() == "시편":
            self.label_4.setText("편")
        else:
            self.label_4.setText("장") 

        jang = [50, 40, 27, 36, 34, 24, 21, 4, 31, 24, 22, 25, 29, 36, 10, 13, 10, 42, 150, 
        31, 12, 8, 66, 52, 5, 48, 12, 14, 3, 9, 1, 4, 7, 3, 3, 3, 2, 14, 4, 28, 16, 24, 21, 28, 16, 16,
        13, 6, 6, 4, 4, 5, 3, 6, 4, 3, 1, 13, 5, 5, 3, 5, 1, 1, 1, 22]
        
        yakja = ["창", "출", "레", "민", "신", "수", "삿", "룻", "삼상", "삼하", "왕상", "왕하", "대상", "대하",
        "스", "느", "에", "욥", "시", "잠", "전", "아", "사", "렘", "애", "겔", "단", "호", "욜", "암", 
        "옵", "욘", "미", "나", "합", "습", "학", "슥", "말", "마", "막", "눅", "요", "행", "롬", "고전", "고후",
        "갈", "엡", "빌", "골", "살전", "살후", "딤전", "딤후", "딛", "몬", "히", "약", "벧전", "벧후", 
        "요1", "요2", "요3", "유", "계"]

        for i in range(jang[num]):
            self.comboBox_3.addItem(str(i+1))
          
    def comboBox_3TextChange(self):
        
        num = self.comboBox.currentIndex()
        self.comboBox_4.clear()
        self.comboBox_6.clear()
        yakja = ["창", "출", "레", "민", "신", "수", "삿", "룻", "삼상", "삼하", "왕상", "왕하", "대상", "대하",
        "스", "느", "에", "욥", "시", "잠", "전", "아", "사", "렘", "애", "겔", "단", "호", "욜", "암", 
        "옵", "욘", "미", "나", "합", "습", "학", "슥", "말", "마", "막", "눅", "요", "행", "롬", "고전", "고후",
        "갈", "엡", "빌", "골", "살전", "살후", "딤전", "딤후", "딛", "몬", "히", "약", "벧전", "벧후", 
        "요일", "요이", "요삼", "유", "계"]

        whereisFile = str("./text/" + str(num + 1) + ".txt")
        with open(whereisFile, 'r') as f:
            lines = f.readlines()
            start = False
            startmemo = 0
            endmemo = 0

            for i in range(len(lines)):
                if lines[i].startswith(yakja[num] + self.comboBox_3.currentText() + ":1 "):
                    startmemo = i+1
                    start = True
                else:
                    if start == True and (lines[i].startswith(yakja[num] + str(int(self.comboBox_3.currentText()) + 1) + ":1 ") or i == len(lines) - 1):
                        endmemo = i
                        start = False
                        if i == len(lines) - 1:
                            jul = endmemo - startmemo + 2
                        else:
                            jul = endmemo - startmemo + 1
                        for i in range(jul):
                            self.comboBox_4.addItem(str(i+1))
                            self.comboBox_6.addItem(str(i+1))
                        break
    
    def makePPT(self):
        prs = Presentation("./template.pptx")
        howmuchtomake = int(int(self.comboBox_6.currentText()) - int(self.comboBox_4.currentText()))
        if howmuchtomake == 0 or howmuchtomake % 2 != 0:
            howmuchtomake += 1
        else:
            howmuchtomake = int(int(howmuchtomake) / 2 + 1)
        
        yakja = ["창", "출", "레", "민", "신", "수", "삿", "룻", "삼상", "삼하", "왕상", "왕하", "대상", "대하",
        "스", "느", "에", "욥", "시", "잠", "전", "아", "사", "렘", "애", "겔", "단", "호", "욜", "암", 
        "옵", "욘", "미", "나", "합", "습", "학", "슥", "말", "마", "막", "눅", "요", "행", "롬", "고전", "고후",
        "갈", "엡", "빌", "골", "살전", "살후", "딤전", "딤후", "딛", "몬", "히", "약", "벧전", "벧후", 
        "요일", "요이", "요삼", "유", "계"]

        count = int(self.comboBox_4.currentText())
        if self.comboBox_4.currentText() > self.comboBox_6.currentText():
            QMessageBox.about(self, "Bible2PPT", "입력한 절을 확인해주세요.")
            

        elif self.comboBox_4.currentText() == self.comboBox_6.currentText(): #1절만
            slide = prs.slides.add_slide(prs.slide_layouts[0]) # 슬라이드 추가
            title = slide.placeholders[10]
            title.text = yakja[self.comboBox.currentIndex()] + " " + self.comboBox_3.currentText() + ":" + self.comboBox_4.currentText()

            jul = slide.placeholders[12]
            num = self.comboBox.currentIndex()
            filename = str("./text/" + str(int(num) + 1) + ".txt")
            with open(filename, 'r') as f:
                lines = f.readlines()
                for i in range(len(lines)):
                    if lines[i].startswith(yakja[num]+self.comboBox_3.currentText()+":"+self.comboBox_4.currentText()+" "):
                        input_text_range = lines[i].split()
                        input_text = ""

                        for text in range(len(input_text_range)):
                            if input_text_range[text].endswith(">"):
                                end_comment = text
                                break
                            else:
                                print(input_text_range[text])
                                end_comment = 0
                        for text in range(end_comment+1, len(input_text_range)):            
                            input_text += input_text_range[text] + " "

            jul.text = input_text
            number = slide.placeholders[11]
            number.text = self.comboBox_4.currentText()
            try:
                filetitle = yakja[num]+" "+self.comboBox_3.currentText()+"장 "+self.comboBox_4.currentText()+"절"
                filename = QFileDialog.getSaveFileName(self, 'Save As PowerPoint File', filetitle, "PowerPoint 프레젠테이션 (*.pptx)")[0]
                try:
                    prs.save(filename)

                except PermissionError:
                    QMessageBox.about(self, "Bible2PPT", "권한이 없습니다. 같은 이름의 파일이 이미 실행 중인지 확인하세요.")
            except:
                print('error')
        else:
            count = int(self.comboBox_4.currentText())
            for i in range(howmuchtomake + 1):
                if count <= int(self.comboBox_6.currentText()):
                    slide = prs.slides.add_slide(prs.slide_layouts[0]) # 슬라이드 추가
                    title = slide.placeholders[10]    
                    title.text = yakja[self.comboBox.currentIndex()] + " " + self.comboBox_3.currentText() + ":" + self.comboBox_4.currentText() + "-" + self.comboBox_6.currentText()
                    number1 = slide.placeholders[11]
                    number1.text = str(count)
                    text1 = slide.placeholders[12]

                    num = self.comboBox.currentIndex()
                    filename = str("./text/" + str(int(num) + 1) + ".txt")
                    with open(filename, 'r') as f:
                        lines = f.readlines()
                        for i in range(len(lines)):
                            if lines[i].startswith(yakja[num]+self.comboBox_3.currentText()+":"+str(count)+" "):
                                input_text_range = lines[i].split()
                                input_text = ""

                        for text in range(len(input_text_range)):
                            if input_text_range[text].endswith(">"):
                                end_comment = text
                                break
                            else:
                                print(input_text_range[text])
                                end_comment = 0
                        for text in range(end_comment+1, len(input_text_range)):            
                            input_text += input_text_range[text] + " "

                    text1.text = input_text
                    count += 1
                    if len(input_text) > 100:
                        next_slide = True
                    else:
                        next_slide = False

                if next_slide == False and count <= int(self.comboBox_6.currentText()):
                    number2 = slide.placeholders[13]
                    

                    text2 = slide.placeholders[14]
                    num = self.comboBox.currentIndex()
                    filename = str("./text/" + str(int(num) + 1) + ".txt")
                    with open(filename, 'r') as f:
                        lines = f.readlines()
                        for i in range(len(lines)):
                            if lines[i].startswith(yakja[num]+self.comboBox_3.currentText()+":"+str(count)+" "):
                                input_text_range = lines[i].split()
                                input_text = ""

                        for text in range(len(input_text_range)):
                            if input_text_range[text].endswith(">"):
                                end_comment = text
                                break
                            else:
                                print(input_text_range[text])
                                end_comment = 0
                        for text in range(end_comment+1, len(input_text_range)):            
                            input_text += input_text_range[text] + " "

                    if len(input_text) < 120:
                        number2.text = str(count)
                        text2.text = input_text
                        count += 1

                if count > int(self.comboBox_6.currentText()):
                    break
                
            try:
                filetitle = yakja[num]+" "+self.comboBox_3.currentText()+"장 "+self.comboBox_4.currentText()+"-"+self.comboBox_6.currentText()+"절"
                filename = QFileDialog.getSaveFileName(self, 'Save As PowerPoint File', filetitle, "PowerPoint 프레젠테이션 (*.pptx)")[0]
                try:
                    prs.save(filename)

                except PermissionError:
                    QMessageBox.about(self, "Bible2PPT", "권한이 없습니다. 같은 이름의 파일이 이미 실행 중인지 확인하세요.")
            except:
                print('error')

if __name__ == "__main__" :
    app = QApplication(sys.argv) 
    myWindow = WindowClass() 
    myWindow.show()
    app.exec_()
